import requests
import json
import re
import plugins
from bridge.reply import Reply, ReplyType
from bridge.context import ContextType
from channel.chat_message import ChatMessage
from plugins import *
from common.log import logger
from common.expired_dict import ExpiredDict
import os
from docx import Document
import markdown
import fitz
from openpyxl import load_workbook
import csv
from bs4 import BeautifulSoup
from pptx import Presentation
from PIL import Image
import base64
import html
import importlib.util
import urllib3

# 禁用不安全请求警告
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# 检查是否安装了openai
try:
    from openai import OpenAI
    has_openai = True
except ImportError:
    has_openai = False

# Moved remove_markdown function here
def remove_markdown(text):
    # 替换Markdown的粗体标记
    text = text.replace("**", "")
    # 替换Markdown的标题标记
    text = text.replace("### ", "").replace("## ", "").replace("# ", "")
    return text

EXTENSION_TO_TYPE = {
    'pdf': 'pdf',
    'doc': 'docx', 'docx': 'docx',
    'md': 'md',
    'txt': 'txt',
    'xls': 'excel', 'xlsx': 'excel',
    'csv': 'csv',
    'html': 'html', 'htm': 'html',
    'ppt': 'ppt', 'pptx': 'ppt'
}

@plugins.register(
    name="sum4all",
    desire_priority=2,
    desc="A plugin for summarizing all things",
    version="0.7.11",
    author="fatwang2",
)

class sum4all(Plugin):
    def __init__(self):
        super().__init__()
        try:
            curdir = os.path.dirname(__file__)
            config_path = os.path.join(curdir, "config.json")
            if os.path.exists(config_path):
                with open(config_path, "r", encoding="utf-8") as f:
                    self.config = json.load(f)
            else:
                # 使用父类的方法来加载配置
                self.config = super().load_config()

                if not self.config:
                    raise Exception("config.json not found")
            # 设置事件处理函数
            self.handlers[Event.ON_HANDLE_CONTEXT] = self.on_handle_context
            self.params_cache = ExpiredDict(300)

            
            # 从配置中提取所需的设置
            self.keys = self.config.get("keys", {})
            self.url_sum = self.config.get("url_sum", {})
            self.search_sum = self.config.get("search_sum", {})
            self.file_sum = self.config.get("file_sum", {})
            self.image_sum = self.config.get("image_sum", {})
            self.note = self.config.get("note", {})

            self.sum4all_key = self.keys.get("sum4all_key", "")
            self.search1api_key = self.keys.get("search1api_key", "")
            self.gemini_key = self.keys.get("gemini_key", "")
            self.bibigpt_key = self.keys.get("bibigpt_key", "")
            self.outputLanguage = self.keys.get("outputLanguage", "zh-CN")
            self.opensum_key = self.keys.get("opensum_key", "")
            self.open_ai_api_key = self.keys.get("open_ai_api_key", "")
            self.model = self.keys.get("model", "gpt-3.5-turbo")
            self.open_ai_api_base = self.keys.get("open_ai_api_base", "https://api.openai.com/v1")
            self.azure_deployment_id = self.keys.get("azure_deployment_id", "")
            self.xunfei_app_id = self.keys.get("xunfei_app_id", "")
            self.xunfei_api_key = self.keys.get("xunfei_api_key", "")
            self.xunfei_api_secret = self.keys.get("xunfei_api_secret", "")
            self.perplexity_key = self.keys.get("perplexity_key", "")
            self.flomo_key = self.keys.get("flomo_key", "")
            # Load Aliyun config
            self.aliyun_key = self.keys.get("aliyun_key", "")
            self.aliyun_base_url = self.keys.get("aliyun_base_url", "https://dashscope.aliyuncs.com/compatible-mode/v1")
            self.aliyun_model = self.keys.get("aliyun_model", "qwen-max")
            self.aliyun_vl_model = self.keys.get("aliyun_vl_model", "qwen-vl-max-latest")
            self.aliyun_sum_model = self.keys.get("aliyun_sum_model", "qwen-long")
            # Load SiliconFlow (sflow) config - CORRECTED
            # self.sflow_key = self.keys.get("sflow_key", "") # Old incorrect key
            # self.sflow_base_url = self.keys.get("sflow_base_url", "") # Old incorrect key
            # self.sflow_model = self.keys.get("sflow_model", "") # Old incorrect key - missing distinction
            self.siliconflow_key = self.keys.get("siliconflow_key", "") # Use correct key from user
            self.siliconflow_base_url = self.keys.get("siliconflow_base_url", "") # Use correct key from user
            self.siliconflow_vl_model = self.keys.get("siliconflow_vl_model", "") # Load VL model
            self.siliconflow_sum_model = self.keys.get("siliconflow_sum_model", "") # Load SUM model

            # 提取sum服务的配置
            self.url_sum_enabled = self.url_sum.get("enabled", False)
            self.url_sum_service = self.url_sum.get("service", "")
            self.url_sum_group = self.url_sum.get("group", True)
            self.url_sum_qa_enabled = self.url_sum.get("qa_enabled", True)
            self.url_sum_qa_prefix = self.url_sum.get("qa_prefix", "问")
            self.url_sum_prompt = self.url_sum.get("prompt", "")

            self.search_sum_enabled = self.search_sum.get("enabled", False)
            self.search_sum_service = self.search_sum.get("service", "")
            self.search_service = self.search_sum.get("search_service", "duckduckgo")
            self.search_sum_group = self.search_sum.get("group", True)
            self.search_sum_search_prefix = self.search_sum.get("search_prefix", "搜")
            self.search_sum_prompt = self.search_sum.get("prompt", "")

            self.file_sum_enabled = self.file_sum.get("enabled", False)
            self.file_sum_service = self.file_sum.get("service", "")
            self.max_file_size = self.file_sum.get("max_file_size", 15000)
            self.file_sum_group = self.file_sum.get("group", True)
            self.file_sum_qa_prefix = self.file_sum.get("qa_prefix", "问")
            self.file_sum_prompt = self.file_sum.get("prompt", "")

            self.image_sum_enabled = self.image_sum.get("enabled", False)
            self.image_sum_service = self.image_sum.get("service", "")
            self.image_sum_group = self.image_sum.get("group", True)
            self.image_sum_qa_prefix = self.image_sum.get("qa_prefix", "问")
            self.image_sum_prompt = self.image_sum.get("prompt", "")

            self.note_enabled = self.note.get("enabled", False)
            self.note_service = self.note.get("service", "")
            self.note_prefix = self.note.get("prefix", "记")

            # 初始化成功日志
            logger.info("[sum4all] inited.")
            # Log configured services after initialization
            if self.url_sum_enabled:
                logger.info(f"[sum4all] URL Summarization ENABLED. Service: {self.url_sum_service}")
            else:
                logger.info("[sum4all] URL Summarization DISABLED.")
            
            if self.file_sum_enabled:
                logger.info(f"[sum4all] File Summarization ENABLED. Service: {self.file_sum_service}")
            else:
                logger.info("[sum4all] File Summarization DISABLED.")

            if self.image_sum_enabled:
                logger.info(f"[sum4all] Image Summarization ENABLED. Service: {self.image_sum_service}")
            else:
                logger.info("[sum4all] Image Summarization DISABLED.")
                
            if self.search_sum_enabled:
                 logger.info(f"[sum4all] Search Summarization ENABLED. Service: {self.search_sum_service}, Search Engine: {self.search_service}")
            else:
                 logger.info("[sum4all] Search Summarization DISABLED.")
                 
            if self.note_enabled:
                 logger.info(f"[sum4all] Note Taking ENABLED. Service: {self.note_service}")
            else:
                 logger.info("[sum4all] Note Taking DISABLED.")
                 
        except Exception as e:
            # 初始化失败日志
            logger.warn(f"sum4all init failed: {e}")
    def on_handle_context(self, e_context: EventContext):
        context = e_context["context"]
        if context.type not in [ContextType.TEXT, ContextType.SHARING,ContextType.FILE,ContextType.IMAGE]:
            return
        msg: ChatMessage = e_context["context"]["msg"]
        user_id = msg.from_user_id
        content = context.content
        isgroup = e_context["context"].get("isgroup", False)

        url_match = re.match('https?://(?:[-\w.]|(?:%[\da-fA-F]{2}))+', content)
        unsupported_urls = re.search(r'.*finder\.video\.qq\.com.*|.*support\.weixin\.qq\.com/update.*|.*support\.weixin\.qq\.com/security.*|.*mp\.weixin\.qq\.com/mp/waerrpage.*', content)

            # 检查输入是否以"搜索前缀词" 开头
        if content.startswith(self.search_sum_search_prefix) and self.search_sum_enabled:
            # 如果消息来自一个群聊，并且你不希望在群聊中启用搜索功能，直接返回
            if isgroup and not self.search_sum_group:
                return
            # Call new function to handle search operation
            self.call_service(content, e_context, "search")
            return
        
        if user_id in self.params_cache and ('last_file_content' in self.params_cache[user_id] or 'last_image_base64' in self.params_cache[user_id] or 'last_url' in self.params_cache[user_id]):
            # 如果存在最近一次处理的文件路径，触发文件理解函数
            if 'last_file_content' in self.params_cache[user_id] and content.startswith(self.file_sum_qa_prefix):
                logger.info('Content starts with the file_sum_qa_prefix.')
                # 去除关键词和紧随其后的空格
                new_content = content[len(self.file_sum_qa_prefix):]
                self.params_cache[user_id]['prompt'] = new_content
                logger.info('params_cache for user has been successfully updated.')            
                self.handle_file(self.params_cache[user_id]['last_file_content'], e_context)
            # 如果存在最近一次处理的图片路径，触发图片理解函数
            elif 'last_image_base64' in self.params_cache[user_id] and content.startswith(self.image_sum_qa_prefix):
                logger.info('Content starts with the image_sum_qa_prefix.')
                # 去除关键词和紧随其后的空格
                new_content = content[len(self.image_sum_qa_prefix):]
                self.params_cache[user_id]['prompt'] = new_content
                logger.info('params_cache for user has been successfully updated.')            
                self.handle_image(self.params_cache[user_id]['last_image_base64'], e_context)

            # 如果存在最近一次处理的URL，触发URL理解函数
            elif 'last_url' in self.params_cache[user_id] and content.startswith(self.url_sum_qa_prefix):
                logger.info('Content starts with the url_sum_qa_prefix.')
                # 去除关键词和紧随其后的空格
                new_content = content[len(self.url_sum_qa_prefix):]
                self.params_cache[user_id]['prompt'] = new_content
                logger.info('params_cache for user has been successfully updated.')            
                self.call_service(self.params_cache[user_id]['last_url'], e_context ,"sum")
            elif 'last_url' in self.params_cache[user_id] and content.startswith(self.note_prefix) and self.note_enabled and not isgroup:
                logger.info('Content starts with the note_prefix.')
                new_content = content[len(self.note_prefix):]
                self.params_cache[user_id]['note'] = new_content
                logger.info('params_cache for user has been successfully updated.')  
                self.call_service(self.params_cache[user_id]['last_url'], e_context, "note")
        if context.type == ContextType.FILE:
            if isgroup and not self.file_sum_group:
                # 群聊中忽略处理文件
                logger.info("群聊消息，文件处理功能已禁用")
                return
            logger.info("on_handle_context: 处理上下文开始")
            context.get("msg").prepare()
            file_path = context.content
            logger.info(f"on_handle_context: 获取到文件路径 {file_path}")
            
            # 检查是否应该进行文件总结
            if self.file_sum_enabled:
                # 更新params_cache中的last_file_content
                self.params_cache[user_id] = {}
                file_content = self.extract_content(file_path)
                if file_content is None:
                    logger.info("文件内容无法提取，跳过处理")
                else:
                    self.params_cache[user_id]['last_file_content'] = file_content
                    logger.info('Updated last_file_content in params_cache for user.')
                    self.handle_file(file_content, e_context)
            else:
                logger.info("文件总结功能已禁用，不对文件内容进行处理")
            # 删除文件
            os.remove(file_path)
            logger.info(f"文件 {file_path} 已删除")
        elif context.type == ContextType.IMAGE:
            if isgroup and not self.image_sum_group:
                # 群聊中忽略处理图片
                logger.info("群聊消息，图片处理功能已禁用")
                return
            logger.info("on_handle_context: 开始处理图片")
            context.get("msg").prepare()
            image_path = context.content
            logger.info(f"on_handle_context: 获取到图片路径 {image_path}")
            
            
            # 检查是否应该进行图片总结
            if self.image_sum_enabled:
                # 将图片路径转换为Base64编码的字符串
                base64_image = self.encode_image_to_base64(image_path)
                # 更新params_cache中的last_image_path
                self.params_cache[user_id] = {}
                self.params_cache[user_id]['last_image_base64'] = base64_image
                logger.info('Updated last_image_base64 in params_cache for user.')
                self.handle_image(base64_image, e_context)

            else:
                logger.info("图片总结功能已禁用，不对图片内容进行处理")
            # 删除文件
            os.remove(image_path)
            logger.info(f"文件 {image_path} 已删除")
        elif context.type == ContextType.SHARING and self.url_sum_enabled:  #匹配卡片分享
            content = html.unescape(content)
            if unsupported_urls:  #匹配不支持总结的卡片
                if isgroup:  ##群聊中忽略
                    return
                else:  ##私聊回复不支持
                    logger.info("[sum4all] Unsupported URL : %s", content)
                    reply = Reply(type=ReplyType.TEXT, content="不支持总结小程序和视频号")
                    e_context["reply"] = reply
                    e_context.action = EventAction.BREAK_PASS
            else:  #匹配支持总结的卡片
                if isgroup:  #处理群聊总结
                    if self.url_sum_group:  #group_sharing = True进行总结，False则忽略。
                        logger.info("[sum4all] Summary URL : %s", content)
                        # 更新params_cache中的last_url
                        self.params_cache[user_id] = {}
                        self.params_cache[user_id]['last_url'] = content
                        logger.info('Updated last_url in params_cache for user.')
                        # Add log before calling call_service
                        logger.info(f"[sum4all] Detected URL/Sharing (Group), preparing to call service for sum. Service configured: {self.url_sum_service}")
                        self.call_service(content, e_context, "sum")
                        return
                    else:
                        return
                else:  #处理私聊总结
                    logger.info("[sum4all] Summary URL : %s", content)
                    # 更新params_cache中的last_url
                    self.params_cache[user_id] = {}
                    self.params_cache[user_id]['last_url'] = content
                    logger.info('Updated last_url in params_cache for user.')
                    # Add log before calling call_service
                    logger.info(f"[sum4all] Detected URL/Sharing (Private), preparing to call service for sum. Service configured: {self.url_sum_service}")
                    self.call_service(content, e_context, "sum")
                    return
            
        elif url_match and self.url_sum_enabled: #匹配URL链接
            if unsupported_urls:  #匹配不支持总结的网址
                logger.info("[sum4all] Unsupported URL : %s", content)
                reply = Reply(type=ReplyType.TEXT, content="不支持总结小程序和视频号")
                e_context["reply"] = reply
                e_context.action = EventAction.BREAK_PASS
            else:
                logger.info("[sum4all] Summary URL : %s", content)
                # 更新params_cache中的last_url
                self.params_cache[user_id] = {}
                self.params_cache[user_id]['last_url'] = content
                logger.info('Updated last_url in params_cache for user.')
                # Add log before calling call_service
                logger.info(f"[sum4all] Detected URL Match, preparing to call service for sum. Service configured: {self.url_sum_service}")
                self.call_service(content, e_context, "sum")
                return
    def call_service(self, content, e_context, service_type):
        # Add log at the beginning of call_service
        logger.info(f"[sum4all] Entering call_service for type '{service_type}'. Configured URL service: {self.url_sum_service}")
        if service_type == "search":
            if self.search_sum_service == "openai" or self.search_sum_service == "sum4all" or self.search_sum_service == "gemini" or self.search_sum_service == "azure":
                self.handle_search(content, e_context)
            elif self.search_sum_service == "perplexity":
                self.handle_perplexity(content, e_context)
        elif service_type == "sum":
            if self.url_sum_service == "bibigpt":
                self.handle_bibigpt(content, e_context)
            elif self.url_sum_service == "openai":
                self.handle_url(content, e_context)
            elif self.url_sum_service == "sum4all":
                self.handle_sum4all(content, e_context)
            elif self.url_sum_service == "gemini":
                self.handle_gemini(content, e_context)
            elif self.url_sum_service == "azure":
                self.handle_azure(content, e_context)
            elif self.url_sum_service == "opensum":
                self.handle_opensum(content, e_context)
            elif self.url_sum_service == "aliyun":
                self.handle_aliyun_url(content, e_context)
            elif self.url_sum_service == "sflow": # Existing elif for sflow
                self.handle_sflow_url(content, e_context) # Call the sflow handler
        elif service_type == "note":
            if self.note_service == "flomo":
                self.handle_note(content, e_context)
    def handle_note(self,link,e_context):
        msg: ChatMessage = e_context["context"]["msg"]
        user_id = msg.from_user_id
        title = self.params_cache[user_id].get('title', '')
        content = self.params_cache[user_id].get('content', '')
        note = self.params_cache[user_id].get('note', '')
        # 将这些内容按照一定的格式整合到一起
        note_content = f"#sum4all\n{title}\n📒笔记：{note}\n{content}\n{link}"
        payload = {"content": note_content}
        # 将这个字典转换为JSON格式
        payload_json = json.dumps(payload)
        # 创建一个POST请求
        url = self.flomo_key
        headers = {'Content-Type': 'application/json'}
        # 发送这个POST请求
        response = requests.post(url, headers=headers, data=payload_json)
        reply = Reply()
        reply.type = ReplyType.TEXT
        if response.status_code == 200 and response.json()['code'] == 0:
            reply.content = f"已发送到{self.note_service}"        
        else:
            reply.content = "发送失败，错误码：" + str(response.status_code)
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS   
    def short_url(self, long_url):
        url = "https://short.fatwang2.com"
        payload = {
            "url": long_url
        }        
        headers = {'Content-Type': "application/json"}
        response = requests.request("POST", url, json=payload, headers=headers)
        if response.status_code == 200:
            res_data = response.json()
            # 直接从返回的 JSON 中获取短链接
            short_url = res_data.get('shorturl', None)  
            
            if short_url:
                return short_url
        return None
    def get_webpage_content(self, url):
        """获取网页内容"""
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
                'Accept-Language': 'zh-CN,zh;q=0.9,en;q=0.8',
                'Accept-Encoding': 'gzip, deflate, br',
                'Connection': 'keep-alive',
                'Upgrade-Insecure-Requests': '1',
                'Cache-Control': 'max-age=0'
            }

            # 添加调试日志
            logger.debug(f"Checking URL for weixin: url='{url}', contains_weixin={'weixin.qq.com' in url}")
            # 修改判断标准：检查URL是否包含 weixin.qq.com
            if "weixin.qq.com" in url:
                logger.info(f"检测到微信相关域名，尝试直接获取内容: {url}")
                # 更新请求头，加入 Referer
                headers['Referer'] = 'https://mp.weixin.qq.com/' # Referer 保持 mp.
                response = requests.get(url, headers=headers, verify=False, timeout=15)
                response.raise_for_status()
                # 记录响应头中的 Content-Type
                content_type_header = response.headers.get('Content-Type')
                logger.debug(f"Response Content-Type header: {content_type_header}")
                # 移除 apparent_encoding 的猜测，强制使用 UTF-8
                # response.encoding = response.apparent_encoding 
                html_content_bytes = response.content
                html_text = html_content_bytes.decode('utf-8', errors='replace')
                # 记录解码后的repr
                logger.debug(f"Decoded html_text (repr): {repr(html_text[:500])}...")
                
                # 使用强制解码后的文本进行解析
                soup = BeautifulSoup(html_text, 'html.parser')
                
                # 查找标题 (可选, 主要用于调试)
                title_tag = soup.find('h1', class_='rich_media_title') or \
                            soup.find('h1', id='activity-name')
                if title_tag:
                    logger.info(f"微信文章标题: {title_tag.get_text(strip=True)}")
                    extracted_title = title_tag.get_text(strip=True)
                else:
                    logger.warning("未找到微信文章标题标签")
                    extracted_title = None

                # 查找正文，优先 rich_media_content
                content_tag = soup.find('div', class_='rich_media_content') or \
                             soup.find('div', id='js_content')
                             
                if content_tag:
                    # 清理不需要的标签
                    for element in content_tag(['script', 'style', 'iframe', 'img', 'video']):
                        element.decompose()
                    # 提取清理后的文本
                    text = content_tag.get_text(separator='\n', strip=True)
                    logger.info("微信公众号内容获取并清理成功")
                    # 记录 get_text 后的 repr
                    logger.debug(f"Text after get_text (repr): {repr(text[:500])}...")
                else:
                    logger.error("无法从微信公众号页面提取正文内容 (rich_media_content 或 js_content)")
                    return None, None # 返回 None 内容和 None 标题
            else:
                # 使用jina.ai预处理URL
                jina_url = f"https://r.jina.ai/{url}"
                logger.info(f"非微信域名，使用jina.ai预处理URL: {jina_url}")
                
                # 获取jina.ai处理后的内容
                response = requests.get(jina_url, headers=headers, verify=False, timeout=20) # 增加超时
                response.raise_for_status()
                # Jina 返回的是纯文本，直接使用
                text = response.text
                logger.info("jina.ai 内容获取成功")

            # --- 通用文本清理逻辑 ---
            # 1. 移除URL链接 (保留微信的清理逻辑，但对Jina可能不需要)
            if "weixin.qq.com" not in url:
                 text = re.sub(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\\\(\\\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+', '', text)

            # 2. 移除邮箱地址
            text = re.sub(r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\\.[a-zA-Z]{2,}', '', text)

            # 3. 移除多余的空格和换行 (保留，对两者都有用)
            text = re.sub(r'\\s+', ' ', text)
            text = text.replace('\\n', ' ') # 将换行符替换为空格，避免过多空行

            # 4. 移除特殊字符 (保留，但对中文内容可能过于激进，稍作调整)
            # text = re.sub(r'[^\w\s\u4e00-\u9fff.,!?，。！？]', '', text) # 原逻辑
            # 修正Linter错误：
            text = re.sub(r'[^\w\s\u4e00-\u9fff,.!?;:"，。！？；："()"（）《》【】「」￥$@%#&*_=+`~^<>|\/\[\]{}-]', '', text) # 统一引号并确保括号闭合

            # 5. 移除数字编号 (保留)
            text = re.sub(r'^\d+\.\s*', '', text, flags=re.MULTILINE)

            # 6. 移除重复的标点符号 (保留)
            text = re.sub(r'([.,!?，。！？])\\1+', r'\\1', text)

            # 7. 移除多余空行 (调整逻辑，先合并空格再处理)
            text = re.sub(r' {2,}', ' ', text) # 合并多个空格
            text = text.strip() # 移除首尾空格

            # 8. 移除行首行尾的空白 (已通过 strip 处理)

            # 如果内容太短，可能是没有正确获取
            if len(text) < 50:
                logger.warning("获取到的内容太短，可能未正确获取文章内容")
                return None, None # 返回 None 内容和 None 标题
            
            # 记录最终返回前的 repr
            logger.debug(f"Final text before return (repr): {repr(text[:500])}...")
            # 如果是 Jina 路径，尝试提取标题
            if "weixin.qq.com" not in url:
                lines = text.split('\n')
                extracted_title = next((line.strip() for line in lines if line.strip()), None)
            # 否则 extracted_title 已在微信路径中设置
            logger.debug(f"Extracted title before return: {extracted_title}")
            return text, extracted_title
                
        except Exception as e:
            logger.error(f"获取网页内容时出错: {e}")
            return None, None # 异常也返回 None, None

    def handle_url(self, content, e_context):
        logger.info('Handling OpenAI request...')
        # 只处理 OpenAI 服务
        if self.url_sum_service != "openai":
            logger.error(f"当前配置的服务不是 OpenAI: {self.url_sum_service}")
            return
            
        api_key = self.open_ai_api_key
        # 修改API基础URL的格式，确保使用http而不是https
        api_base = self.open_ai_api_base.replace('https://', 'http://')
        model = self.model
        
        msg: ChatMessage = e_context["context"]["msg"]
        user_id = msg.from_user_id
        user_params = self.params_cache.get(user_id, {})
        prompt = user_params.get('prompt', self.url_sum_prompt)
        
        # 获取网页内容
        webpage_content = self.get_webpage_content(content)
        if not webpage_content:
            reply_content = "无法获取网页内容，请检查链接是否有效"
        else:
            # 构建 OpenAI API 请求
            headers = {
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {api_key}'
            }
            
            # 构建系统提示词
            system_prompt = """你是一个专业的网页内容总结专家。请按照以下格式总结网页内容：
1. 首先用一句话总结文章的核心观点（30字以内）
2. 然后列出3-5个关键要点
3. 使用emoji让表达更生动
4. 保持专业、客观的语气"""

            # 构建用户提示词
            user_prompt = f"""请总结以下网页内容：
{prompt}

网页内容：
{webpage_content[:4000]}  # 限制内容长度，避免超出token限制"""

            # 构建请求体
            payload = {
                "model": model,
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                "temperature": 0.7,
                "max_tokens": 1000
            }

            additional_content = ""
            try:
                logger.info('Sending request to OpenAI...')
                # 直接调用 OpenAI API，禁用SSL验证
                response = requests.post(
                    f"{api_base}/chat/completions",
                    headers=headers,
                    json=payload,
                    verify=False  # 禁用SSL验证
                )
                response.raise_for_status()
                logger.info('Received response from OpenAI.')
                
                response_data = response.json()
                if "choices" in response_data and len(response_data["choices"]) > 0:
                    content = response_data["choices"][0]["message"]["content"]
                    self.params_cache[user_id]['content'] = content
                    
                    # 尝试从内容中提取标题（第一行）
                    lines = content.split('\n')
                    if lines:
                        title = lines[0].strip()
                        self.params_cache[user_id]['title'] = title
                        if title:
                            additional_content += f"{title}\n\n"
                    
                    reply_content = additional_content + content
                else:
                    reply_content = "无法获取有效的响应内容"

            except requests.exceptions.RequestException as e:
                logger.error(f"Error calling OpenAI API: {e}")
                reply_content = f"调用 OpenAI API 时发生错误: {str(e)}"

        reply = Reply()
        reply.type = ReplyType.TEXT
        if not self.url_sum_qa_enabled:
            reply.content = remove_markdown(reply_content)
        elif isgroup or not self.note_enabled:
            reply.content = f"{remove_markdown(reply_content)}\n\n💬5min内输入{self.url_sum_qa_prefix}+问题，可继续追问"
        elif self.note_enabled:
            reply.content = f"{remove_markdown(reply_content)}\n\n💬5min内输入{self.url_sum_qa_prefix}+问题，可继续追问\n💡输入{self.note_prefix}+笔记，可保存到{self.note_service}"
        
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    def handle_bibigpt(self, content, e_context):    
        headers = {
            'Content-Type': 'application/json'
        }
        payload_params = {
            "url": content,
            "includeDetail": False,
            "promptConfig": {
                "outputLanguage": self.outputLanguage
            }
        }

        payload = json.dumps(payload_params)           
        try:
            api_url = f"https://bibigpt.co/api/open/{self.bibigpt_key}"
            response = requests.request("POST",api_url, headers=headers, data=payload)
            response.raise_for_status()
            data = json.loads(response.text)
            summary_original = data.get('summary', 'Summary not available')
            html_url = data.get('htmlUrl', 'HTML URL not available')
            # 获取短链接
            short_url = self.short_url(html_url) 
            
            # 如果获取短链接失败，使用 html_url
            if short_url is None:
                short_url = html_url if html_url != 'HTML URL not available' else 'URL not available'
            
            # 移除 "##摘要"、"## 亮点" 和 "-"
            summary = summary_original.split("详细版（支持对话追问）")[0].replace("## 摘要\n", "📌总结：").replace("## 亮点\n", "").replace("- ", "")
        except requests.exceptions.RequestException as e:
            reply = f"An error occurred"

        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{summary}\n\n详细链接：{short_url}"

        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    def handle_opensum(self, content, e_context):
        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {self.opensum_key}'
        }
        payload = json.dumps({"link": content})
        try:
            api_url = "https://read.thinkwx.com/api/v1/article/summary"
            response = requests.request("POST",api_url, headers=headers, data=payload)
            response.raise_for_status()
            data = json.loads(response.text)
            summary_data = data.get('data', {})  # 获取data字段                
            summary_original = summary_data.get('summary', 'Summary not available')
            # 使用正则表达式提取URL
            url_pattern = r'https:\/\/[^\s]+'
            match = re.search(url_pattern, summary_original)
            html_url = match.group(0) if match else 'HTML URL not available'            
            # 获取短链接
            short_url = self.short_url(html_url) if match else html_url
            # 用于移除摘要中的URL及其后的所有内容
            url_pattern_remove = r'https:\/\/[^\s]+[\s\S]*'
            summary = re.sub(url_pattern_remove, '', summary_original).strip()        

        except requests.exceptions.RequestException as e:
            summary = f"An error occurred"
            short_url = 'URL not available'
        
        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{summary}\n\n详细链接：{short_url}"

        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS    
    def handle_search(self, content, e_context):
        # 根据sum_service的值选择API密钥和基础URL
        if self.search_sum_service == "openai":
            api_key = self.open_ai_api_key
            api_base = self.open_ai_api_base
            model = self.model
            headers = {
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {api_key}'
            }
            
            # 构建系统提示词
            system_prompt = """你是一个专业的搜索专家。请根据用户的问题，从搜索结果中提取相关信息，并按照以下格式回答：
1. 首先用一句话总结答案的核心观点（30字以内）
2. 然后列出3-5个关键要点
3. 使用emoji让表达更生动
4. 保持专业、客观的语气"""

            # 构建用户提示词
            user_prompt = f"""请根据以下问题搜索并总结：
{content[len(self.search_sum_search_prefix):]}

搜索服务：{self.search_service}"""

            # 构建请求体
            payload = {
                "model": model,
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                "temperature": 0.7,
                "max_tokens": 1000
            }

            try:
                response = requests.post(
                    f"{api_base}/chat/completions",
                    headers=headers,
                    json=payload
                )
                response.raise_for_status()
                response_data = response.json()
                
                if "choices" in response_data and len(response_data["choices"]) > 0:
                    reply_content = response_data["choices"][0]["message"]["content"]
                else:
                    reply_content = "无法获取有效的响应内容"

            except requests.exceptions.RequestException as e:
                logger.error(f"Error calling OpenAI API: {e}")
                reply_content = f"调用 OpenAI API 时发生错误: {str(e)}"

        elif self.search_sum_service == "sum4all":
            api_key = self.sum4all_key
            api_base = "https://pro.sum4all.site/v1"
            model = "sum4all"
            headers = {
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {api_key}'
            }
            payload = {
                "ur": content[len(self.search_sum_search_prefix):],
                "prompt": self.search_sum_prompt,
                "model": model,
                "base": api_base,
                "search1api_key": self.search1api_key,
                "search_service": self.search_service
            }
            try:
                response = requests.post(api_base, headers=headers, json=payload)
                response.raise_for_status()
                response_data = response.json()
                if response_data.get("success"):
                    reply_content = response_data["content"].replace("\\n", "\n")
                else:
                    reply_content = "无法获取有效的响应内容"
            except requests.exceptions.RequestException as e:
                logger.error(f"Error calling Sum4All API: {e}")
                reply_content = f"调用 Sum4All API 时发生错误: {str(e)}"

        elif self.search_sum_service == "gemini":
            api_key = self.gemini_key
            model = "gemini"
            api_base = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent"
            headers = {
                'Content-Type': 'application/json',
                'x-goog-api-key': api_key
            }
            
            system_prompt = """你是一个专业的搜索专家。请根据用户的问题，从搜索结果中提取相关信息，并按照以下格式回答：
1. 首先用一句话总结答案的核心观点（30字以内）
2. 然后列出3-5个关键要点
3. 使用emoji让表达更生动
4. 保持专业、客观的语气"""

            user_prompt = f"""请根据以下问题搜索并总结：
{content[len(self.search_sum_search_prefix):]}

搜索服务：{self.search_service}"""

            payload = {
                "contents": [
                    {"role": "user", "parts": [{"text": system_prompt}]},
                    {"role": "model", "parts": [{"text": "okay"}]},
                    {"role": "user", "parts": [{"text": user_prompt}]}
                ],
                "generationConfig": {
                    "maxOutputTokens": 800
                }
            }

            try:
                response = requests.post(api_base, headers=headers, json=payload)
                response.raise_for_status()
                response_data = response.json()
                
                if "candidates" in response_data and len(response_data["candidates"]) > 0:
                    reply_content = response_data["candidates"][0]["content"]["parts"][0]["text"]
                else:
                    reply_content = "无法获取有效的响应内容"

            except requests.exceptions.RequestException as e:
                logger.error(f"Error calling Gemini API: {e}")
                reply_content = f"调用 Gemini API 时发生错误: {str(e)}"

        elif self.search_sum_service == "azure":
            api_key = self.open_ai_api_key
            api_base = f"{self.open_ai_api_base}/openai/deployments/{self.azure_deployment_id}/chat/completions?api-version=2024-02-15-preview"
            model = self.model
            headers = {
                'Content-Type': 'application/json',
                'api-key': api_key
            }
            
            system_prompt = """你是一个专业的搜索专家。请根据用户的问题，从搜索结果中提取相关信息，并按照以下格式回答：
1. 首先用一句话总结答案的核心观点（30字以内）
2. 然后列出3-5个关键要点
3. 使用emoji让表达更生动
4. 保持专业、客观的语气"""

            user_prompt = f"""请根据以下问题搜索并总结：
{content[len(self.search_sum_search_prefix):]}

搜索服务：{self.search_service}"""

            payload = {
                "model": model,
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ]
            }

            try:
                response = requests.post(api_base, headers=headers, json=payload)
                response.raise_for_status()
                response_data = response.json()
                
                if "choices" in response_data and len(response_data["choices"]) > 0:
                    reply_content = response_data["choices"][0]["message"]["content"]
                else:
                    reply_content = "无法获取有效的响应内容"

            except requests.exceptions.RequestException as e:
                logger.error(f"Error calling Azure API: {e}")
                reply_content = f"调用 Azure API 时发生错误: {str(e)}"

        else:
            logger.error(f"未知的search_service配置: {self.search_sum_service}")
            return

        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{remove_markdown(reply_content)}"            
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    def handle_perplexity(self, content, e_context):

        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {self.perplexity_key}'
        }
        data = {
            "model": "sonar-small-online",
            "messages": [
                {"role": "system", "content": self.search_sum_prompt},
                {"role": "user", "content": content}
        ]
        }
        try:
            api_url = "https://api.perplexity.ai/chat/completions"
            response = requests.post(api_url, headers=headers, json=data)
            response.raise_for_status()
            # 处理响应数据
            response_data = response.json()
            # 这里可以根据你的需要处理响应数据
            # 解析 JSON 并获取 content
            if "choices" in response_data and len(response_data["choices"]) > 0:
                first_choice = response_data["choices"][0]
                if "message" in first_choice and "content" in first_choice["message"]:
                    content = first_choice["message"]["content"]
                else:
                    print("Content not found in the response")
            else:
                print("No choices available in the response")
        except requests.exceptions.RequestException as e:
            # 处理可能出现的错误
            logger.error(f"Error calling perplexity: {e}")
        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{remove_markdown(content)}"            
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    def get_help_text(self, verbose=False, **kwargs):
        help_text = "Help you summarize all things\n"
        if not verbose:
            return help_text
        help_text += "1.Share me the link and I will summarize it for you\n"
        help_text += f"2.{self.search_sum_search_prefix}+query,I will search online for you\n"
        return help_text
    def handle_file(self, content, e_context):
        logger.info("handle_file: 向LLM发送内容总结请求")
        # 根据sum_service的值选择API密钥和基础URL
        if self.file_sum_service == "openai":
            api_key = self.open_ai_api_key
            api_base = self.open_ai_api_base
            model = self.model
        elif self.file_sum_service == "azure":
            api_key = self.open_ai_api_key
            api_base = f"{self.open_ai_api_base}/openai/deployments/{self.azure_deployment_id}/chat/completions?api-version=2024-02-15-preview"
            model = self.model
        elif self.file_sum_service == "sum4all":
            api_key = self.sum4all_key
            api_base = "https://pro.sum4all.site/v1"
            model = "sum4all"
        elif self.file_sum_service == "gemini":
            api_key = self.gemini_key
            model = "gemini"
            api_base = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent"
        elif self.file_sum_service == "aliyun":
            reply_content = self.handle_aliyun_file(content, e_context)
            
            reply = Reply()
            reply.type = ReplyType.TEXT
            reply.content = f"{remove_markdown(reply_content)}\n\n💬5min内输入{self.file_sum_qa_prefix}+问题，可继续追问" 
            e_context["reply"] = reply
            e_context.action = EventAction.BREAK_PASS
            return
        else:
            logger.error(f"未知的sum_service配置: {self.file_sum_service}")
            return
        msg: ChatMessage = e_context["context"]["msg"]
        user_id = msg.from_user_id
        user_params = self.params_cache.get(user_id, {})
        prompt = user_params.get('prompt', self.file_sum_prompt)
        if model == "gemini":
            headers = {
                'Content-Type': 'application/json',
                'x-goog-api-key': api_key
            }
            data = {
            "contents": [
                {"role": "user", "parts": [{"text": prompt}]},
                {"role": "model", "parts": [{"text": "okay"}]},
                {"role": "user", "parts": [{"text": content}]}
            ],
            "generationConfig": {
                "maxOutputTokens": 800
            }
            }
            api_url = api_base
        elif self.file_sum_service == "azure":
            headers = {
                "Content-Type": "application/json",
                "api-key": api_key
            }
            data = {
                "model": model,
                "messages": [
                    {"role": "system", "content": prompt},
                    {"role": "user", "content": content}
                ]
            }
            api_url = api_base
        elif self.file_sum_service == "aliyun":
            api_key = self.aliyun_key
            model = "aliyun"
            api_base = self.aliyun_base_url
            
            if has_openai:
                # 使用OpenAI客户端库
                try:
                    logger.info(f"使用OpenAI客户端调用阿里云API: {api_base}")
                    client = OpenAI(
                        api_key=api_key,
                        base_url=api_base
                    )
                    
                    completion = client.chat.completions.create(
                        model=self.aliyun_sum_model,
                        messages=[
                            {"role": "system", "content": prompt},
                            {"role": "user", "content": content}
                        ],
                        temperature=0.7,
                        max_tokens=2000
                    )
                    
                    logger.info("OpenAI客户端成功获取响应")
                    response_content = completion.choices[0].message.content.strip()
                    return response_content.replace("\\n", "\n")
                    
                except Exception as e:
                    logger.error(f"使用OpenAI客户端调用阿里云API出错: {e}")
                    logger.info("转为使用requests直接调用")
            
            # 使用requests直接调用
            try:
                logger.info("使用requests直接调用阿里云API")
                headers = {
                    'Content-Type': 'application/json',
                    'Authorization': f'Bearer {api_key}'
                }
                
                data = {
                    "model": self.aliyun_sum_model,
                    "messages": [
                        {"role": "system", "content": prompt},
                        {"role": "user", "content": content}
                    ],
                    "temperature": 0.7,
                    "max_tokens": 2000
                }
                
                api_url = api_base if "/chat/completions" in api_base else f"{api_base}/chat/completions"
                logger.info(f"请求URL: {api_url}")
                
                response = requests.post(
                    api_url,
                    headers=headers,
                    json=data,
                    verify=False,
                    timeout=30
                )
                
                response.raise_for_status()
                logger.info(f"API响应状态码: {response.status_code}")
                
                response_data = response.json()
                if "choices" in response_data and len(response_data["choices"]) > 0:
                    first_choice = response_data["choices"][0]
                    if "message" in first_choice and "content" in first_choice["message"]:
                        response_content = first_choice["message"]["content"].strip()
                        logger.info("成功获取阿里云API响应内容")
                        return response_content.replace("\\n", "\n")
                    else:
                        logger.error("阿里云API响应中未找到内容字段")
                        return "未能从阿里云API获取有效的响应内容"
                else:
                    logger.error("阿里云API响应中未找到choices字段")
                    return "未能从阿里云API获取有效的响应内容"
                    
            except Exception as e:
                logger.error(f"调用阿里云API时出错: {e}")
                if hasattr(e, 'response') and e.response:
                    logger.error(f"响应状态码: {e.response.status_code}")
                    logger.error(f"响应内容: {e.response.text}")
                return f"调用阿里云API时发生错误: {str(e)}"
        else:
            headers = {
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {api_key}'
            }
            data = {
                "model": model,
                "messages": [
                    {"role": "system", "content": prompt},
                    {"role": "user", "content": content}
                ]
            }
            api_url = f"{api_base}/chat/completions"
        try:
            response = requests.post(api_url, headers=headers, data=json.dumps(data))
            response.raise_for_status()
            response_data = response.json()
            
            # 解析 JSON 并获取 content
            if model == "gemini":
                if "candidates" in response_data and len(response_data["candidates"]) > 0:
                    first_candidate = response_data["candidates"][0]
                    if "content" in first_candidate:
                        if "parts" in first_candidate["content"] and len(first_candidate["content"]["parts"]) > 0:
                            response_content = first_candidate["content"]["parts"][0]["text"].strip()  # 获取响应内容
                            logger.info(f"Gemini API response content: {response_content}")  # 记录响应内容
                            reply_content = response_content.replace("\\n", "\n")  # 替换 \\n 为 \n
                        else:
                            logger.error("Parts not found in the Gemini API response content")
                            reply_content = "Parts not found in the Gemini API response content"
                    else:
                        logger.error("Content not found in the Gemini API response candidate")
                        reply_content = "Content not found in the Gemini API response candidate"
                else:
                    logger.error("No candidates available in the Gemini API response")
                    reply_content = "No candidates available in the Gemini API response"        
            else:
                if "choices" in response_data and len(response_data["choices"]) > 0:
                    first_choice = response_data["choices"][0]
                    if "message" in first_choice and "content" in first_choice["message"]:
                        response_content = first_choice["message"]["content"].strip()
                        logger.info(f"LLM API response content")  # 记录响应内容
                        reply_content = response_content.replace("\\n", "\n")  # 替换 \\n 为 \n
                    else:
                        logger.error("Content not found in the response")
                        reply_content = "Content not found in the LLM API response"
                else:
                    logger.error("No choices available in the response")
                    reply_content = "No choices available in the LLM API response"

        except requests.exceptions.RequestException as e:
            logger.error(f"Error calling LLM API: {e}")
            reply_content = f"An error occurred while calling LLM API"

        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{remove_markdown(reply_content)}\n\n💬5min内输入{self.file_sum_qa_prefix}+问题，可继续追问" 
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    def read_pdf(self, file_path):
        logger.info(f"开始读取PDF文件：{file_path}")
        doc = fitz.open(file_path)
        content = ' '.join([page.get_text() for page in doc])
        logger.info(f"PDF文件读取完成：{file_path}")

        return content
    def read_word(self, file_path):
        doc = Document(file_path)
        return ' '.join([p.text for p in doc.paragraphs])
    def read_markdown(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            md_content = file.read()
            return markdown.markdown(md_content)
    def read_excel(self, file_path):
        workbook = load_workbook(file_path)
        content = ''
        for sheet in workbook:
            for row in sheet.iter_rows():
                content += ' '.join([str(cell.value) for cell in row])
                content += '\n'
        return content
    def read_txt(self, file_path):
        logger.debug(f"开始读取TXT文件: {file_path}")
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            logger.debug(f"TXT文件读取完成: {file_path}")
            logger.debug("TXT文件内容的前50个字符：")
            logger.debug(content[:50])  # 打印文件内容的前50个字符
            return content
        except Exception as e:
            logger.error(f"读取TXT文件时出错: {file_path}，错误信息: {str(e)}")
            return ""
    def read_csv(self, file_path):
        content = ''
        with open(file_path, 'r', encoding='utf-8') as csvfile:
            reader = csv.reader(csvfile)
            for row in reader:
                content += ' '.join(row) + '\n'
        return content
    def read_html(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file, 'html.parser')
            return soup.get_text()
    def read_ppt(self, file_path):
        presentation = Presentation(file_path)
        content = ''
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + '\n'
        return content
    def extract_content(self, file_path):
        logger.info(f"extract_content: 提取文件内容，文件路径: {file_path}")
        file_size = os.path.getsize(file_path) // 1000  # 将文件大小转换为KB
        if file_size > int(self.max_file_size):
            logger.warning(f"文件大小超过限制({self.max_file_size}KB),不进行处理。文件大小: {file_size}KB")
            return None
        file_extension = os.path.splitext(file_path)[1][1:].lower()
        logger.info(f"extract_content: 文件类型为 {file_extension}")

        file_type = EXTENSION_TO_TYPE.get(file_extension)

        if not file_type:
            logger.error(f"不支持的文件扩展名: {file_extension}")
            return None

        read_func = {
            'pdf': self.read_pdf,
            'docx': self.read_word,
            'md': self.read_markdown,
            'txt': self.read_txt,
            'excel': self.read_excel,
            'csv': self.read_csv,
            'html': self.read_html,
            'ppt': self.read_ppt
        }.get(file_type)

        if not read_func:
            logger.error(f"不支持的文件类型: {file_type}")
            return None
        logger.info("extract_content: 文件内容提取完成")
        return read_func(file_path)
    def encode_image_to_base64(self, image_path):
        # 打开图片
        img = Image.open(image_path)
        # 只有当图片的宽度大于1024像素时，才调整图片大小
        if img.width > 1024:
            img = img.resize((1024, int(img.height*1024/img.width)))
            # 将调整大小后的图片保存回原文件
            img.save(image_path)

        # 打开调整大小后的图片，读取并进行base64编码
        with open(image_path, "rb") as image_file:
            encoded = base64.b64encode(image_file.read()).decode('utf-8')
        return encoded
    def handle_image(self, base64_image, e_context):
        logger.info("handle_image: 解析图像处理API的响应")
        msg: ChatMessage = e_context["context"]["msg"]
        user_id = msg.from_user_id
        user_params = self.params_cache.get(user_id, {})
        prompt = user_params.get('prompt', self.image_sum_prompt)

        logger.debug(f"[sum4all] Checking image_sum_service value: '{self.image_sum_service}'") # ADDED DEBUG LOG
        if self.image_sum_service == "openai":
            api_key = self.open_ai_api_key
            api_base = f"{self.open_ai_api_base}/chat/completions"
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {api_key}"
            }
            model = "gpt-4o-mini"
        elif self.image_sum_service == "azure":
            api_key = self.open_ai_api_key
            api_base = f"{self.open_ai_api_base}/openai/deployments/{self.azure_deployment_id}/chat/completions?api-version=2024-02-15-preview"
            headers = {
                "Content-Type": "application/json",
                "api-key": api_key
            }
            model = "gpt-4o-mini"
        elif self.image_sum_service == "xunfei":
            api_key = self.xunfei_api_key
            api_base = "https://spark.sum4all.site/v1/chat/completions"
            model = "spark-chat-vision"
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {api_key}"
            }
        elif self.image_sum_service == "sum4all":
            api_key = self.sum4all_key
            api_base = "https://pro.sum4all.site/v1/chat/completions"
            model = "sum4all-vision"
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {api_key}"
            }
        elif self.image_sum_service == "gemini":
            api_key = self.gemini_key
            api_base = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent"
            payload = {
                "contents": [
                    {
                        "parts": [
                            {"text": prompt},
                            {
                                "inline_data": {
                                    "mime_type":"image/png",
                                    "data": base64_image
                                }
                            }
                        ]
                    }
                ],
                "generationConfig": {
                    "maxOutputTokens": 800
                }
            }
            headers = {
                "Content-Type": "application/json",
                "x-goog-api-key": api_key
            }
        elif self.image_sum_service == "aliyun":
            api_key = self.aliyun_key
            api_base = self.aliyun_base_url
            
            if has_openai:
                # 使用OpenAI客户端库
                try:
                    client = OpenAI(
                        api_key=api_key,
                        base_url=api_base
                    )
                    
                    completion = client.chat.completions.create(
                        model=self.aliyun_vl_model,
                        messages=[
                            {"role": "system", "content": self.image_sum_prompt},
                            {"role": "user", "content": [{"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}]}
                        ],
                        temperature=0.7,
                        max_tokens=2000
                    )
                    
                    reply_content = completion.choices[0].message.content
                    
                except Exception as e:
                    logger.error(f"Error using OpenAI client for Aliyun API: {e}")
                    # 失败后回退到直接使用requests
                    headers = {
                        'Content-Type': 'application/json',
                        'Authorization': f'Bearer {api_key}'
                    }
                    payload = {
                        "model": self.aliyun_vl_model,
                        "messages": [
                            {"role": "system", "content": self.image_sum_prompt},
                            {"role": "user", "content": [{"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}]}
                        ],
                        "temperature": 0.7,
                        "max_tokens": 2000
                    }
                    response = requests.post(
                        api_base if "/chat/completions" in api_base else f"{api_base}/chat/completions",
                        headers=headers,
                        json=payload,
                        verify=False,
                        timeout=30
                    )
                    response.raise_for_status()
                    result = response.json()
                    if "choices" in result and len(result["choices"]) > 0:
                        reply_content = result["choices"][0]["message"]["content"]
                    else:
                        logger.error("阿里百炼 API 返回格式错误")
                        reply_content = "总结失败，请稍后重试"
            # The following else block is the fallback for Aliyun if has_openai is false or fails
            else: # This else corresponds to 'if has_openai' within the aliyun block
                logger.info("[sum4all][aliyun_image] Using requests directly for Aliyun...")
                headers = {
                    'Content-Type': 'application/json',
                    'Authorization': f'Bearer {api_key}'
                }
                payload = {
                    "model": self.aliyun_vl_model,
                    "messages": [
                        {"role": "system", "content": self.image_sum_prompt},
                        {"role": "user", "content": [{"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}]}
                    ],
                    "temperature": 0.7,
                    "max_tokens": 2000
                }
                try:
                    response = requests.post(
                        api_base if "/chat/completions" in api_base else f"{api_base}/chat/completions",
                        headers=headers,
                        json=payload,
                        verify=False,
                        timeout=30
                    )
                    response.raise_for_status()
                    result = response.json()
                    if "choices" in result and len(result["choices"]) > 0:
                        reply_content = result["choices"][0]["message"]["content"]
                    else:
                        logger.error("[sum4all][aliyun_image] Aliyun API response format error (requests).")
                        reply_content = "总结失败，请稍后重试 (requests)"
                except requests.exceptions.RequestException as e:
                    logger.error(f"[sum4all][aliyun_image] Error calling Aliyun API via requests: {e}")
                    reply_content = "调用阿里百炼API失败 (requests)"

        # <<<<< CORRECTED SFLOW BLOCK STARTS HERE >>>>>
        elif self.image_sum_service == "sflow":
            logger.info("[sum4all][sflow_image] Handling SiliconFlow image request...")
            api_key = self.siliconflow_key
            api_base = self.siliconflow_base_url
            model = self.siliconflow_vl_model # Use the specific VL model

            if not api_key or not api_base or not model:
                 logger.error("[sum4all][sflow_image] SiliconFlow configuration (key, base_url, or vl_model) is missing.")
                 reply_content = "SiliconFlow 配置不完整，无法处理图片"
            else:
                headers = {
                    "Content-Type": "application/json",
                    "Authorization": f"Bearer {api_key}"
                }
                payload = {
                    "model": model,
                    "messages": [
                        {
                            "role": "user",
                            "content": [
                                {"type": "text", "text": prompt},
                                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                            ]
                        }
                    ],
                    "temperature": 0.7,
                    "max_tokens": 1500
                }
                logger.debug(f"[sum4all][sflow_image] Sending payload: {payload}")
                reply_content = "处理图片时发生未知错误" # Default error
                try:
                    logger.info(f"[sum4all][sflow_image] Sending request to {api_base}/chat/completions")
                    response = requests.post(
                        f"{api_base}/chat/completions",
                        headers=headers,
                        json=payload,
                        verify=False,
                        timeout=60
                    )
                    logger.info(f"[sum4all][sflow_image] Received response status: {response.status_code}")
                    response.raise_for_status()
                    response_data = response.json()
                    logger.debug(f"[sum4all][sflow_image] Response data: {response_data}")
                    if "choices" in response_data and len(response_data["choices"]) > 0:
                        result = response_data["choices"][0].get("message", {}).get("content")
                        if result:
                            reply_content = result
                            logger.info("[sum4all][sflow_image] Successfully extracted content from response.")
                        else:
                            logger.error("[sum4all][sflow_image] 'content' missing in response choice.")
                            reply_content = "无法从 SiliconFlow API 获取有效的响应内容 (content missing)"
                    else:
                        logger.error("[sum4all][sflow_image] 'choices' missing or empty in response.")
                        reply_content = "无法从 SiliconFlow API 获取有效的响应内容 (choices missing)"
                except requests.exceptions.Timeout:
                    logger.error("[sum4all][sflow_image] Request timed out.")
                    reply_content = "调用 SiliconFlow API 处理图片超时"
                except requests.exceptions.RequestException as e:
                    logger.error(f"[sum4all][sflow_image] API request error: {e}")
                    if hasattr(e, 'response') and e.response is not None:
                        logger.error(f"[sum4all][sflow_image] Response status code: {e.response.status_code}")
                        logger.error(f"[sum4all][sflow_image] Response content: {e.response.text}")
                        reply_content = f"调用 SiliconFlow API 处理图片出错: Status {e.response.status_code}"
                    else:
                        reply_content = f"调用 SiliconFlow API 处理图片时发生网络错误: {str(e)}"

        # <<<<< END OF CORRECTED SFLOW BLOCK >>>>>
        else: # Final else for the main service selection
            logger.error(f"未知的image_sum_service配置: {self.image_sum_service}")
            return

        reply = Reply()
        reply.type = ReplyType.TEXT
        reply.content = f"{remove_markdown(reply_content)}\n\n💬5min内输入{self.image_sum_qa_prefix}+问题，可继续追问"
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS
    
    def handle_sum4all(self, content, e_context):
        logger.info('Handling Sum4All request...')
        # 由于sum4all.site服务不可用，我们改用OpenAI API
        api_key = self.open_ai_api_key
        api_base = self.open_ai_api_base
        model = self.model
        
        msg: ChatMessage = e_context["context"]["msg"]
        user_id = msg.from_user_id
        user_params = self.params_cache.get(user_id, {})
        isgroup = e_context["context"].get("isgroup", False)
        prompt = user_params.get('prompt', self.url_sum_prompt)
        
        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {api_key}'
        }
        
        # 构建系统提示词
        system_prompt = """你是一个专业的网页内容总结专家。请按照以下格式总结网页内容：
1. 首先用一句话总结文章的核心观点（30字以内）
2. 然后列出3-5个关键要点
3. 使用emoji让表达更生动
4. 保持专业、客观的语气"""

        # 构建用户提示词
        user_prompt = f"""请总结以下网页内容：
{prompt}

网页链接：{content}"""

        # 构建请求体
        payload = {
            "model": model,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            "temperature": 0.7,
            "max_tokens": 1000
        }

        try:
            response = requests.post(
                f"{api_base}/chat/completions",
                headers=headers,
                json=payload,
                verify=False  # 禁用SSL验证
            )
            response.raise_for_status()
            response_data = response.json()
            
            if "choices" in response_data and len(response_data["choices"]) > 0:
                content = response_data["choices"][0]["message"]["content"]
                self.params_cache[user_id]['content'] = content
                
                # 尝试从内容中提取标题（第一行）
                lines = content.split('\n')
                if lines:
                    title = lines[0].strip()
                    self.params_cache[user_id]['title'] = title
                
                additional_content = ""
                if title:
                    additional_content += f"{title}\n\n"
                reply_content = additional_content + content
            else:
                reply_content = "无法获取有效的响应内容"

        except requests.exceptions.RequestException as e:
            logger.error(f"Error calling OpenAI API: {e}")
            reply_content = f"调用 OpenAI API 时发生错误: {str(e)}"

        reply = Reply()
        reply.type = ReplyType.TEXT
        if not self.url_sum_qa_enabled:
            reply.content = remove_markdown(reply_content)
        elif isgroup or not self.note_enabled:
            reply.content = f"{remove_markdown(reply_content)}\n\n💬5min内输入{self.url_sum_qa_prefix}+问题，可继续追问"
        elif self.note_enabled:
            reply.content = f"{remove_markdown(reply_content)}\n\n💬5min内输入{self.url_sum_qa_prefix}+问题，可继续追问\n💡输入{self.note_prefix}+笔记，可保存到{self.note_service}"
        
        e_context["reply"] = reply
        e_context.action = EventAction.BREAK_PASS

    def handle_gemini(self, content, e_context):
        logger.info('Handling Gemini request...')
        # 获取网页内容
        webpage_content = self.get_webpage_content(content)
        if not webpage_content:
            reply_content = "无法获取网页内容，请检查链接是否有效"
        else:
            api_key = self.gemini_key
            model = "gemini"
            api_base = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent"
            
            msg: ChatMessage = e_context["context"]["msg"]
            user_id = msg.from_user_id
            user_params = self.params_cache.get(user_id, {})
            isgroup = e_context["context"].get("isgroup", False)
            prompt = user_params.get('prompt', self.url_sum_prompt)
            
            headers = {
                'Content-Type': 'application/json',
                'x-goog-api-key': api_key
            }
            
            # 构建系统提示词
            system_prompt = """你是一个专业的网页内容总结专家。请按照以下格式总结网页内容：
1. 首先用一句话总结文章的核心观点（30字以内）
2. 然后列出3-5个关键要点
3. 使用emoji让表达更生动
4. 保持专业、客观的语气"""
